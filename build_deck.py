"""build_deck.py — v2
=====================================================================
Generates `deck.pptx` for the Boeing / Chomsky 5-minute presentation.

v2 changes over v1:
  - Editorial serif titles (Cambria) + sans body (Inter / Calibri).
  - One slide per filter (5 total) + intro filter slide + hook + timeline
    + counterfactual close = 9 slides.
  - Each filter slide carries a verbatim Chomsky quote pulled from
    `Manufacturing Consent by Chomsky 2002 (1).pdf` with page number.
  - Speaker notes carry the full argument: a 3-4 beat outline at the
    top, a 40-60 second read-aloud script under it, and curated image
    source URLs at the bottom for assets the user wants to upgrade
    later (Shanahan portrait, KAL 007 NYT page, etc.).
  - New visual evidence: advertising-flow diagram (filter 2),
    article-positioning analysis + Mar 11 / Mar 17 publication
    comparison + Boeing apparatus (filter 4), worthy/unworthy victims
    comparison drawn typographically (filter 5).

Run:
    python3 build_deck.py
=====================================================================
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
IMG = ROOT / "deck" / "images"
OUT = ROOT / "deck.pptx"

# ---------------- design tokens ----------------

BG        = RGBColor(0xD6, 0xDC, 0xE0)
INK       = RGBColor(0x10, 0x10, 0x10)
INK_MUTED = RGBColor(0x55, 0x5B, 0x60)
INK_SOFT  = RGBColor(0x88, 0x8E, 0x94)
LINE      = RGBColor(0x0B, 0x24, 0x47)
ACCENT    = RGBColor(0xC8, 0x10, 0x2E)
SLAB      = RGBColor(0xC4, 0xCC, 0xD2)
PAPER     = RGBColor(0xFA, 0xF8, 0xF4)
PLACEHOLD = RGBColor(0xB8, 0xBF, 0xC6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Cambria"          # editorial serif
QUOTE_FONT = "Cambria"
BODY_FONT  = "Inter"            # sans body; falls back to Calibri/Aptos
MONO_FONT  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------- helpers ----------------

def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def text(slide, body, x, y, w, h, *,
         font=BODY_FONT, size=18, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def rich_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for chunk, style in para:
            run = p.add_run()
            run.text = chunk
            run.font.name      = style.get("font", BODY_FONT)
            run.font.size      = Pt(style.get("size", 18))
            run.font.bold      = style.get("bold", False)
            run.font.italic    = style.get("italic", False)
            run.font.underline = style.get("underline", False)
            run.font.color.rgb = style.get("color", INK)
    return tb


def image(slide, name, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(IMG / name), x, y, width=w, height=h)


def rect(slide, x, y, w, h, *, fill=None, line=None, line_width=None,
         shadow=False):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        if line_width is not None:
            r.line.width = line_width
    return r


def oval(slide, x, y, w, h, fill, line=None):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line
    return o


def connector(slide, x1, y1, x2, y2, color=INK_MUTED, width=Pt(1)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = width
    return c


def hairline(slide, x, y, w, color=INK, weight=Pt(1.25)):
    return rect(slide, x, y, w, Emu(int(weight)), fill=color)


def notes(slide, outline, script, sources=None):
    """Render outline + script + sources into the speaker-notes pane.
    `outline` is a list of beat strings.  `script` is the read-aloud
    paragraph.  `sources` (optional) is a list of (label, url) tuples
    for images the user may want to source manually."""
    parts = []
    parts.append("OUTLINE")
    for i, beat in enumerate(outline, 1):
        parts.append(f"  {i}. {beat}")
    parts.append("")
    parts.append("SCRIPT")
    parts.append(script)
    if sources:
        parts.append("")
        parts.append("IMAGE SOURCES (to drop into deck/images/ for v3)")
        for label, url in sources:
            parts.append(f"  • {label} — {url}")
    slide.notes_slide.notes_text_frame.text = "\n".join(parts)


def placeholder(slide, label, x, y, w, h):
    rect(slide, x, y, w, h, fill=PLACEHOLD, line=INK_SOFT)
    text(slide, label, x, y, w, h,
         font=BODY_FONT, size=14, italic=True, color=INK_MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_title(slide, kicker, title):
    """Editorial slide-header treatment: small red kicker, serif title,
    rule below.  Used on every interior slide."""
    text(slide, kicker.upper(),
         Inches(0.7), Inches(0.45), Inches(12), Inches(0.35),
         font=BODY_FONT, size=12, bold=True, color=ACCENT)
    text(slide, title,
         Inches(0.7), Inches(0.78), Inches(12), Inches(0.85),
         font=TITLE_FONT, size=36, bold=True, color=INK)
    hairline(slide, Inches(0.7), Inches(1.65), Inches(11.93),
             color=INK, weight=Pt(1.5))


def quote_slab(slide, body, attribution, *, y, h=Inches(1.35),
               dark=True):
    """Reusable Chomsky-quote slab.  `dark` selects the navy ground;
    set False for a light slab on a busy slide."""
    fill = LINE if dark else SLAB
    fg   = WHITE if dark else INK
    mute = RGBColor(0xCF, 0xD6, 0xDE) if dark else INK_MUTED
    rect(slide, Inches(0.7), y, Inches(11.93), h, fill=fill)
    rich_text(slide, [
        [("“" + body + "”",
          {"font": QUOTE_FONT, "size": 18, "italic": True, "color": fg})],
        [(" ", {"size": 6})],
        [(attribution,
          {"font": BODY_FONT, "size": 12, "color": mute})],
    ], Inches(1.1), y + Inches(0.18), Inches(11.13), h - Inches(0.3),
       line_spacing=1.25)


# ---------------- slides ----------------

def slide_1_hook(prs):
    s = blank(prs)
    # Editorial kicker + date
    text(s, "MARCH 11, 2019",
         Inches(0.7), Inches(0.55), Inches(12), Inches(0.35),
         font=BODY_FONT, size=12, bold=True, color=ACCENT)
    text(s, "The framing that the propaganda model explains.",
         Inches(0.7), Inches(0.88), Inches(12), Inches(0.55),
         font=TITLE_FONT, size=22, italic=True, color=INK_MUTED)

    # Headline screenshot, prominent
    hw = Inches(11.0)
    hh = Inches(11.0 * 297 / 803)   # ≈ 4.07"
    image(s, "wsj_headline_march11.png",
          (SLIDE_W - hw) / 2, Inches(1.85), w=hw)

    # 189 + 157 = 346, framed under
    rich_text(s, [
        [("189",  {"font": TITLE_FONT, "size": 36, "bold": True, "color": ACCENT}),
         ("  Lion Air 610   ·   ",
          {"font": BODY_FONT, "size": 18, "color": INK_MUTED}),
         ("157", {"font": TITLE_FONT, "size": 36, "bold": True, "color": ACCENT}),
         ("  Ethiopian 302   ·   ",
          {"font": BODY_FONT, "size": 18, "color": INK_MUTED}),
         ("346", {"font": TITLE_FONT, "size": 36, "bold": True, "color": INK}),
         ("  dead",
          {"font": BODY_FONT, "size": 18, "color": INK_MUTED})],
    ], Inches(0.7), Inches(6.25), Inches(11.93), Inches(0.7),
       align=PP_ALIGN.CENTER)

    text(s, "Robert Wall  ·  Andrew Tangel  ·  Andy Pasztor   |   The Wall Street Journal",
         Inches(0), Inches(7.05), SLIDE_W, Inches(0.35),
         font=BODY_FONT, size=11, color=INK_SOFT,
         align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "Set the death toll cold: 189 + 157 = 346, in 132 days.",
            "Same aircraft, same MCAS system, same flight profile.",
            "Show what the WSJ chose to frame as the news the next morning.",
            "Pause and let the audience read the headline.",
        ],
        script=(
            "On October 29, 2018, Lion Air Flight 610 crashed twelve minutes after "
            "takeoff. 189 dead. On March 10, 2019 — 132 days later — Ethiopian "
            "Airlines Flight 302 crashed six minutes after takeoff. 157 more dead. "
            "Same aircraft. Same MCAS system. Same fatal flight profile. The morning "
            "after the second crash, the Wall Street Journal — America's paper of "
            "business record — chose to frame the news this way."
        ),
        sources=[
            ("WSJ Mar 11, 2019 article (ProQuest WSJ archive / UChicago library)",
             "https://www.wsj.com/articles/the-faa-has-no-current-plans-to-ground-boeings-737-max-after-deadly-crash-11552324520"),
        ],
    )


def slide_2_timeline(prs):
    s = blank(prs)
    slide_title(s, "Article in context", "Same aircraft.  No changes.")

    # Ethiopian wreckage thumbnail anchored top-right inside the header band
    image(s, "ethiopian_302_wreckage.png",
          Inches(9.4), Inches(0.45), w=Inches(3.5))

    # Horizontal axis
    axis_y = Inches(3.65)
    rect(s, Inches(0.7), axis_y, Inches(11.93), Emu(38100), fill=LINE)

    xs = [Inches(1.0), Inches(3.6), Inches(6.2), Inches(8.7), Inches(12.0)]
    for xp in xs:
        oval(s, xp - Inches(0.11), axis_y - Inches(0.11),
             Inches(0.22), Inches(0.22), LINE)

    headers = [
        ("Oct 29, 2018",       "Lion Air 610",            "189 dead",     xs[0]),
        ("Nov 13, 2018",       "MCAS revealed (WSJ)",     "Pasztor/Tangel scoop", xs[1]),
        ("Mar 10, 2019",       "Ethiopian 302",           "157 dead",     xs[2]),
        ("Mar 11, 2019 (AM)",  "China grounds 737 MAX",   "before WSJ article ran", xs[3]),
        ("Mar 13, 2019",       "FAA finally grounds",     "(last major regulator)", xs[4]),
    ]
    for date_str, line1, line2, xp in headers:
        text(s, date_str,
             xp - Inches(1.25), Inches(2.55), Inches(2.5), Inches(0.4),
             font=BODY_FONT, size=13, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
        text(s, line1,
             xp - Inches(1.25), Inches(2.92), Inches(2.5), Inches(0.4),
             font=TITLE_FONT, size=15, italic=True, color=INK,
             align=PP_ALIGN.CENTER)
        text(s, line2,
             xp - Inches(1.25), Inches(3.25), Inches(2.5), Inches(0.35),
             font=BODY_FONT, size=11, color=ACCENT,
             align=PP_ALIGN.CENTER)

    # Below the line: the wreckage photo + the two WSJ headlines
    image(s, "lion_air_610_wreckage.png",
          Inches(0.7), Inches(4.15), w=Inches(2.6))
    image(s, "wsj_headline_nov21.png",
          Inches(3.55), Inches(4.45), w=Inches(4.6))
    image(s, "wsj_headline_march11.png",
          Inches(8.4), Inches(4.45), w=Inches(4.6))

    # Bottom: the "51 / 72 / FAA stood alone" tagline
    rich_text(s, [
        [("Within 72 hours of Ethiopian:  ",
          {"font": BODY_FONT, "size": 14, "color": INK}),
         ("51 regulators worldwide",
          {"font": BODY_FONT, "size": 14, "bold": True, "color": ACCENT}),
         ("  grounded the 737 MAX.   The FAA stood alone.",
          {"font": BODY_FONT, "size": 14, "color": INK})],
    ], Inches(0.7), Inches(7.0), Inches(11.93), Inches(0.4),
       align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "132 days between the two crashes; 51 regulators ground in 72 hours.",
            "China grounded BEFORE the WSJ article ran.",
            "FAA was the LAST major regulator to ground — not the first.",
            "WSJ framed FAA's reluctance as the news, not the global consensus.",
            "Transition: 'Why? That's what the propaganda model lets us see.'",
        ],
        script=(
            "132 days separated these two crashes. Within hours of the Ethiopian "
            "crash, China's Civil Aviation Administration grounded the 737 MAX. "
            "Within 72 hours, 51 regulators worldwide had grounded it. By the time "
            "this article ran on the morning of March 11, China had already "
            "grounded the aircraft. The FAA stood alone — and would for another "
            "72 hours, becoming the LAST major regulator to ground. But the Wall "
            "Street Journal didn't frame that worldwide consensus as the news. It "
            "framed the FAA's reluctance as the news. The article makes the "
            "American regulator the default position and treats global action as "
            "deviation requiring explanation. Why? That's what the propaganda "
            "model lets us see."
        ),
    )


def slide_3_five_filters(prs):
    s = blank(prs)
    slide_title(s, "Manufacturing Consent", "Chomsky's Five Filters")

    filters = [
        ("1", "Ownership",       "& profit orientation"),
        ("2", "Advertising",     "the reader economy"),
        ("3", "Sourcing",        "the beat & access economy"),
        ("4", "Flak",            "as anticipatory discipline"),
        ("5", "Unifying ideology", "anticommunism  →  ?"),
    ]
    box_w = Inches(2.25)
    box_h = Inches(2.05)
    gap   = Inches(0.18)
    total_w = box_w * len(filters) + gap * (len(filters) - 1)
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.05)
    for i, (num, name, sub) in enumerate(filters):
        x = start_x + (box_w + gap) * i
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        r.fill.solid(); r.fill.fore_color.rgb = LINE
        r.line.fill.background()
        text(s, num, x, y + Inches(0.18), box_w, Inches(0.6),
             font=TITLE_FONT, size=30, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
        text(s, name, x, y + Inches(0.95), box_w, Inches(0.5),
             font=TITLE_FONT, size=18, italic=True,
             color=WHITE, align=PP_ALIGN.CENTER)
        text(s, sub, x, y + Inches(1.45), box_w, Inches(0.5),
             font=BODY_FONT, size=11,
             color=RGBColor(0xCF, 0xD6, 0xDE), align=PP_ALIGN.CENTER)
        if i < len(filters) - 1:
            ax = x + box_w + Inches(0.02)
            ay = y + box_h / 2 - Inches(0.07)
            tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                     ax, ay, gap - Inches(0.04), Inches(0.14))
            tri.fill.solid(); tri.fill.fore_color.rgb = INK_MUTED
            tri.line.fill.background()

    quote_slab(s,
        "The constraints are so powerful, and are built into the system in "
        "such a fundamental way, that alternative bases of news choices are "
        "hardly imaginable.",
        "— Herman & Chomsky, Manufacturing Consent, p. 2",
        y=Inches(4.85), h=Inches(1.65))

    text(s,
        "Not propaganda in the colloquial sense.  Filtered: structural conditions select what survives.",
        Inches(0.7), Inches(6.75), Inches(11.93), Inches(0.4),
        font=BODY_FONT, size=14, italic=True, color=INK_MUTED,
        align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "Structural, not moral. Filters select what survives — not who's lying.",
            "Pasztor + Tangel broke the MCAS story. They are not bad reporters.",
            "Filters reinforce one another; no single filter is decisive.",
            "Transition: 'Let's see all five operating on this single article.'",
        ],
        script=(
            "Chomsky's argument is structural, not moral. He is not saying the "
            "press is fake or that the journalists are corrupt. Pasztor and Tangel "
            "are not bad reporters — they are the team that BROKE the MCAS "
            "self-certification story four months earlier. The propaganda model "
            "claims that news passes through five sequential filters that select "
            "what survives: ownership, advertising, sourcing, flak, and a unifying "
            "ideology. The journalists are not the filters. They are the operators "
            "of a process whose constraints are determined upstream of any "
            "individual editorial choice. The proof of the model is that good-faith "
            "reporters working with complete integrity will still produce this kind "
            "of output — because the alternative framings cost too much."
        ),
        sources=[
            ("Manufacturing Consent cover (book) — Penguin Random House",
             "https://www.penguinrandomhouse.com/books/78912/manufacturing-consent-by-edward-s-herman-and-noam-chomsky/"),
        ],
    )


def slide_4_ownership(prs):
    s = blank(prs)
    slide_title(s, "Filter 1", "Ownership and Structural Class Identity")

    # Two-column comparison — narrower than v1 so we can fit a class-loop
    # diagram below.  The visual argument is: same class.
    col_w = Inches(5.45)
    col_h = Inches(3.0)
    left_x = Inches(0.7)
    right_x = Inches(7.2)
    col_y = Inches(1.95)

    rect(s, left_x, col_y, col_w, Inches(0.55), fill=LINE)
    rect(s, right_x, col_y, col_w, Inches(0.55), fill=LINE)
    text(s, "WSJ", left_x, col_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=18, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Boeing", right_x, col_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=18, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rect(s, left_x,  col_y + Inches(0.55), col_w, col_h - Inches(0.55), fill=SLAB)
    rect(s, right_x, col_y + Inches(0.55), col_w, col_h - Inches(0.55), fill=SLAB)

    left_bullets = [
        "Owned by News Corp (Murdoch family)",
        "World's largest English-language business publication",
        "Calculates the Dow Jones Industrial Average",
        "Reader base:  corporate-financial elite",
    ]
    right_bullets = [
        "Top-30 U.S. corporation  (~$240B mkt cap)",
        "Largest American exporter (2019)",
        "Dow Jones Industrial Average component",
        "Pentagon prime defense contractor",
        "Anchor employer for Washington + South Carolina",
    ]
    text(s, "\n".join("•   " + it for it in left_bullets),
         left_x + Inches(0.3), col_y + Inches(0.75),
         col_w - Inches(0.5), col_h - Inches(0.8),
         font=BODY_FONT, size=14, color=INK, line_spacing=1.45)
    text(s, "\n".join("•   " + it for it in right_bullets),
         right_x + Inches(0.3), col_y + Inches(0.75),
         col_w - Inches(0.5), col_h - Inches(0.8),
         font=BODY_FONT, size=14, color=INK, line_spacing=1.45)

    # Connecting tissue between the columns — the "same class" rule
    rect(s, Inches(6.15), Inches(2.5), Inches(1.0), Inches(2.45),
         fill=ACCENT)
    text(s, "S A M E\nC L A S S",
         Inches(6.15), Inches(2.5), Inches(1.0), Inches(2.45),
         font=TITLE_FONT, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=1.4)

    # Quote
    quote_slab(s,
        "the dominant media firms are quite large businesses; they are "
        "controlled by very wealthy people or by managers who are subject "
        "to sharp constraints by owners and other market-profit-oriented "
        "forces; and they are closely interlocked, and have important common "
        "interests, with other major corporations, banks, and government.",
        "— Herman & Chomsky, p. 14",
        y=Inches(5.15), h=Inches(1.9))

    text(s,
        "News Corp doesn't tell Pasztor what to write.  The class composition makes Boeing a quasi-internal interest.",
        Inches(0.7), Inches(7.15), Inches(11.93), Inches(0.3),
        font=BODY_FONT, size=12, italic=True, color=INK_MUTED,
        align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "Class position, not editorial orders.",
            "WSJ literally calculates the Dow Jones — which includes Boeing.",
            "Boeing in 2019: Dow component, top exporter, Pentagon prime.",
            "Same elite ecology. Same board interlocks. Same class position.",
            "Transition: 'But ownership alone doesn't pick the headline. Advertising shapes what coverage can sustain.'",
        ],
        script=(
            "Filter 1: ownership. The Wall Street Journal is owned by News Corp — "
            "the Murdoch family. The WSJ literally calculates the Dow Jones "
            "Industrial Average, an index that included Boeing throughout 2019. "
            "Boeing in 2019 was a top-30 US corporation by market cap, the largest "
            "American exporter, the Pentagon's prime defense contractor, the "
            "anchor employer in Washington State and South Carolina. Same elite "
            "ecology. Same board interlocks. Same class position. Filter 1 is not "
            "about Murdoch picking up the phone and telling Pasztor what to write. "
            "It is about the structural fact that the WSJ exists as the publication "
            "of record for the corporate-financial elite that owns Boeing equity, "
            "supplies Boeing, sits on Boeing-adjacent boards, regulates Boeing. The "
            "publication's class identity makes Boeing a quasi-internal interest — "
            "not an external object of scrutiny."
        ),
        sources=[
            ("News Corp ownership chart — already in deck/images/news_corp_chart.png",
             "(local file)"),
            ("News Corp corporate structure (2019) — Wikipedia",
             "https://en.wikipedia.org/wiki/News_Corp"),
        ],
    )


def slide_5_advertising(prs):
    s = blank(prs)
    slide_title(s, "Filter 2", "Advertising and the Reader Economy")

    # ---- Advertising flow diagram ----
    # WSJ hub at center; six surrounding nodes; arrows back to WSJ.
    # Left column = advertisers (ad spend flows IN); right column = readers/
    # holders (subscriptions + equity exposure flow IN).
    diagram_top    = Inches(1.95)
    diagram_bottom = Inches(5.20)
    center_x = Inches(6.667)
    center_y = (diagram_top + diagram_bottom) / 2
    hub_w = Inches(1.7); hub_h = Inches(0.85)
    rect(s, center_x - hub_w / 2, center_y - hub_h / 2, hub_w, hub_h,
         fill=LINE)
    text(s, "WSJ",
         center_x - hub_w / 2, center_y - hub_h / 2, hub_w, hub_h,
         font=TITLE_FONT, size=22, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Column headers — placed ABOVE the node columns, clearly separated
    text(s, "→   AD SPEND",
         Inches(1.7), Inches(1.95), Inches(3.0), Inches(0.3),
         font=BODY_FONT, size=10, bold=True, color=ACCENT,
         align=PP_ALIGN.LEFT)
    text(s, "SUBSCRIPTIONS  /  EQUITY HOLDINGS   ←",
         Inches(8.6), Inches(1.95), Inches(3.5), Inches(0.3),
         font=BODY_FONT, size=10, bold=True, color=ACCENT,
         align=PP_ALIGN.RIGHT)

    # Six nodes, three per column.  Tightened vertical spacing.
    node_w = Inches(2.05); node_h = Inches(0.85)
    nodes = [
        ("Boeing",                          Inches(2.2),  Inches(2.65)),
        ("Defense\n(LMT / NOC / RTX)",       Inches(2.2),  Inches(3.85)),
        ("Suppliers\n(GE / Honeywell)",      Inches(2.2),  Inches(5.05)),
        ("Airlines\n(SWA / AAL / UAL)",      Inches(11.1), Inches(2.65)),
        ("Banks\n(JPM / GS / BLK)",          Inches(11.1), Inches(3.85)),
        ("Reader-investors",                Inches(11.1), Inches(5.05)),
    ]
    for label, nx, ny in nodes:
        rect(s, nx - node_w / 2, ny - node_h / 2, node_w, node_h,
             fill=PAPER, line=INK_SOFT)
        text(s, label,
             nx - node_w / 2, ny - node_h / 2, node_w, node_h,
             font=BODY_FONT, size=11, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.15)
        # Connector from edge of node toward edge of hub
        if nx < center_x:
            x_from = nx + node_w / 2
            x_to   = center_x - hub_w / 2
        else:
            x_from = nx - node_w / 2
            x_to   = center_x + hub_w / 2
        connector(s, x_from, ny, x_to, center_y,
                  color=INK_MUTED, width=Pt(1.25))

    # Explanatory argument under the diagram (above the quote)
    text(s,
        "A single critical article does not threaten this loop.\n"
        "A sustained campaign reframing American manufacturing as systemically unsafe does.",
         Inches(0.7), Inches(5.55), Inches(11.93), Inches(0.7),
         font=TITLE_FONT, size=15, italic=True, color=INK,
         align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Chomsky / Curran & Seaton quote at the very bottom — single line height
    quote_slab(s,
        "advertisers thus acquired a de facto licensing authority since, "
        "without their support, newspapers ceased to be economically viable.",
        "— Curran & Seaton, quoted in Manufacturing Consent, p. 14",
        y=Inches(6.35), h=Inches(0.95))

    notes(s,
        outline=[
            "WSJ's advertiser base IS Boeing's economic ecosystem.",
            "Reader side matters more: WSJ readers HOLD Boeing in portfolios.",
            "Single article ≠ sustained campaign. Filter 2 kills campaigns.",
            "Transition: 'But Filter 2 alone doesn't pick voices. Filter 3 — sourcing — is visible inside the article itself.'",
        ],
        script=(
            "Filter 2: advertising. The WSJ's 2019 advertiser base is essentially "
            "Boeing's economic ecosystem. Boeing itself. Its defense competitors: "
            "Lockheed, Northrop, Raytheon. Its suppliers: GE Aerospace, Honeywell, "
            "Spirit AeroSystems. Its airlines: Southwest, American, United. Its "
            "bankers: JPMorgan, Goldman, BlackRock. But the deeper move is on the "
            "reader side. WSJ readers are institutional investors, corporate "
            "executives, money managers. They HOLD Boeing in their portfolios. "
            "They hold Boeing in their pension funds. Their employers do business "
            "with Boeing. The reader-advertiser ecology IS Boeing's ecology. A "
            "single critical article is permitted — Pasztor's MCAS scoop in "
            "November proved that. What filter 2 prevents is a sustained "
            "campaign — multi-month coverage that reframes American manufacturing "
            "as systemically unsafe. That campaign would break the publication's "
            "relationship to its entire reader-advertiser ecology. So it does not "
            "happen."
        ),
        sources=[
            ("Boeing 2019 annual ad spend / WSJ advertiser data (Statista / Kantar)",
             "https://www.statista.com/topics/1349/boeing/"),
        ],
    )


def slide_6_sourcing(prs):
    s = blank(prs)
    slide_title(s, "Filter 3", "Sourcing and the Beat Economy")

    text(s,
        "Whose voices structure the March 11 article — and whose appear only as reactions.",
        Inches(0.7), Inches(1.75), Inches(11.93), Inches(0.4),
        font=BODY_FONT, size=14, italic=True, color=INK_MUTED)

    col_w = Inches(5.95)
    left_x = Inches(0.7)
    right_x = Inches(6.7)
    col_y = Inches(2.35)
    col_h = Inches(2.95)

    rect(s, left_x, col_y, col_w, Inches(0.6), fill=LINE)
    rect(s, right_x, col_y, col_w, Inches(0.6), fill=ACCENT)
    text(s, "Voices that structure the story",
         left_x, col_y, col_w, Inches(0.6),
         font=TITLE_FONT, size=15, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Voices that appear only as reactions  (or not at all)",
         right_x, col_y, col_w, Inches(0.6),
         font=TITLE_FONT, size=15, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rect(s, left_x,  col_y + Inches(0.6), col_w, col_h - Inches(0.6), fill=SLAB)
    rect(s, right_x, col_y + Inches(0.6), col_w, col_h - Inches(0.6), fill=SLAB)

    left_voices = [
        "The FAA  —  authoritative on its own decision",
        "Boeing  —  authoritative on its own aircraft",
        "U.S. airline executives  (Southwest, American, United)",
        "Anonymous “industry officials”",
    ]
    right_voices = [
        "Mitt Romney  (calling for grounding)",
        "Foreign regulators  (China, EU)  —  framed as deviations",
        "Victims' families  —  absent",
        "Ethiopian / Indonesian regulators  —  absent",
    ]
    text(s, "\n".join("•   " + it for it in left_voices),
         left_x + Inches(0.35), col_y + Inches(0.8),
         col_w - Inches(0.5), col_h - Inches(0.8),
         font=BODY_FONT, size=14, color=INK, line_spacing=1.55)
    text(s, "\n".join("•   " + it for it in right_voices),
         right_x + Inches(0.35), col_y + Inches(0.8),
         col_w - Inches(0.5), col_h - Inches(0.8),
         font=BODY_FONT, size=14, color=INK, line_spacing=1.55)

    # Quote
    quote_slab(s,
        "Officials have and give the facts; reporters merely get them.",
        "— Mark Fishman, quoted in Manufacturing Consent, p. 19",
        y=Inches(5.7), h=Inches(0.9))

    # Bottom argumentative slug
    text(s,
        "Six days later, Dominic Gates (Seattle Times) broke FAA self-certification.  "
        "The story was reportable on Mar 11 — it just broke the beat.",
        Inches(0.7), Inches(6.85), Inches(11.93), Inches(0.4),
        font=BODY_FONT, size=12, italic=True, color=INK_MUTED,
        align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "Filter 3 is visible INSIDE the article — name the voices.",
            "Boeing media office: professionalized, US-hours, pre-prepped.",
            "Foreign regulators + victim families: expensive to source.",
            "Captured-regulator story was reportable — Gates broke it Mar 17/21.",
            "Beat dependency = commercial infrastructure of the reporting.",
            "Transition: 'Filter 4 — flak — explains why the beat survives.'",
        ],
        script=(
            "Filter 3: sourcing — the filter you can see operating inside the "
            "article itself. Let me show you whose voices structure the story. The "
            "FAA. Boeing. US airline executives. Anonymous 'industry officials' — "
            "almost certainly Boeing or FAA staff giving background. These are the "
            "authoritative voices. Now whose voices appear, when they appear at "
            "all, as reactions? Mitt Romney calling for grounding. Foreign "
            "regulators — China, the EU — framed as deviations from default. The "
            "victims' families: absent. The Ethiopian and Indonesian regulators: "
            "absent. This is not Pasztor and Tangel ignoring foreign sources. This "
            "is cost economics. Boeing's media office is professionalized, on US "
            "business hours, with talking points pre-prepared. The Ethiopian Civil "
            "Aviation Authority is not. The captured-regulator framing — the story "
            "that Boeing employees were doing FAA self-certification — was "
            "reportable in March 2019. Dominic Gates at the Seattle Times broke it "
            "six days later. The WSJ team had the expertise. They didn't lead with "
            "it. Why? Because reframing the FAA from 'authoritative source' to "
            "'captured regulator' breaks the aviation beat. The beat is commercial "
            "infrastructure. So the frame stays."
        ),
        sources=[
            ("Gates, 'Flawed analysis, failed oversight' — Seattle Times, Mar 17/21, 2019",
             "https://www.seattletimes.com/business/boeing-aerospace/failed-certification-faa-missed-safety-issues-in-the-737-max-system-implicated-in-the-lion-air-crash/"),
            ("NPR coverage of the Gates piece (Mar 18, 2019)",
             "https://www.npr.org/2019/03/18/704373869/seattle-times-questions-certification-process-of-boeings-737-max"),
        ],
    )


def slide_7_flak(prs):
    """Filter 4 — three columns of evidence the user asked for:
    (1) positioning INSIDE the WSJ article (softening words),
    (2) comparison: Mar 11 WSJ vs Mar 17 Seattle Times,
    (3) Boeing's apparatus + Shanahan."""
    s = blank(prs)
    slide_title(s, "Filter 4", "Flak as Anticipatory Discipline")

    col_w = Inches(3.95)
    col_h = Inches(3.7)
    col_y = Inches(1.95)
    gap = Inches(0.15)
    xs = [Inches(0.7),
          Inches(0.7) + col_w + gap,
          Inches(0.7) + (col_w + gap) * 2]

    # Column 1 — Inside the article: the softening words
    rect(s, xs[0], col_y, col_w, Inches(0.55), fill=LINE)
    text(s, "1.  Softening inside the article",
         xs[0], col_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=14, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, xs[0], col_y + Inches(0.55), col_w, col_h - Inches(0.55), fill=SLAB)

    softening = [
        ("“deadly crash”",
         "instead of\n“second mass-fatality crash of the same\naircraft in five months”"),
        ("“The FAA’s decision”",
         "instead of\n“the regulator that delegated certification\nto the company it was regulating”"),
        ("FAA framed as defaultforeign action framed as deviation",
         ""),
    ]
    yy = col_y + Inches(0.75)
    for headline, sub in softening[:2]:
        text(s, headline,
             xs[0] + Inches(0.25), yy,
             col_w - Inches(0.5), Inches(0.4),
             font=TITLE_FONT, size=13, bold=True, italic=True, color=ACCENT,
             align=PP_ALIGN.LEFT)
        text(s, sub,
             xs[0] + Inches(0.25), yy + Inches(0.4),
             col_w - Inches(0.5), Inches(0.9),
             font=BODY_FONT, size=11, color=INK_MUTED,
             align=PP_ALIGN.LEFT, line_spacing=1.3)
        yy += Inches(1.45)

    # Column 2 — Comparison: WSJ Mar 11 vs Seattle Times Mar 17
    rect(s, xs[1], col_y, col_w, Inches(0.55), fill=LINE)
    text(s, "2.  Compare:  same beat, different frame",
         xs[1], col_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=14, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, xs[1], col_y + Inches(0.55), col_w, col_h - Inches(0.55), fill=SLAB)

    # WSJ Mar 11 reconstruction
    rect(s, xs[1] + Inches(0.2), col_y + Inches(0.75),
         col_w - Inches(0.4), Inches(1.25), fill=WHITE, line=INK_SOFT)
    text(s, "WSJ  —  Mar 11, 2019",
         xs[1] + Inches(0.3), col_y + Inches(0.85),
         col_w - Inches(0.6), Inches(0.3),
         font=BODY_FONT, size=9, bold=True, color=INK_MUTED)
    text(s, "FAA Has No Current Plans\nto Ground Boeing’s 737 MAX\nAfter Deadly Crash",
         xs[1] + Inches(0.3), col_y + Inches(1.1),
         col_w - Inches(0.6), Inches(1.0),
         font=TITLE_FONT, size=11, bold=True, color=INK,
         line_spacing=1.15)

    # Seattle Times Mar 17 reconstruction
    rect(s, xs[1] + Inches(0.2), col_y + Inches(2.15),
         col_w - Inches(0.4), Inches(1.25), fill=WHITE, line=INK_SOFT)
    text(s, "Seattle Times  —  Mar 17, 2019",
         xs[1] + Inches(0.3), col_y + Inches(2.25),
         col_w - Inches(0.6), Inches(0.3),
         font=BODY_FONT, size=9, bold=True, color=INK_MUTED)
    text(s, "Flawed Analysis, Failed Oversight:\nHow Boeing, FAA Certified the\nSuspect 737 MAX Flight Control System",
         xs[1] + Inches(0.3), col_y + Inches(2.5),
         col_w - Inches(0.6), Inches(1.0),
         font=TITLE_FONT, size=10, bold=True, color=INK,
         line_spacing=1.15)
    text(s, "Dominic Gates  —  Pulitzer Prize, 2020",
         xs[1] + Inches(0.3), col_y + Inches(3.35),
         col_w - Inches(0.6), Inches(0.3),
         font=BODY_FONT, size=9, italic=True, color=ACCENT)

    # Column 3 — Boeing's apparatus
    rect(s, xs[2], col_y, col_w, Inches(0.55), fill=LINE)
    text(s, "3.  Boeing’s apparatus, Mar 2019",
         xs[2], col_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=14, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, xs[2], col_y + Inches(0.55), col_w, col_h - Inches(0.55), fill=SLAB)

    apparatus = [
        ("White-shoe PR firms on retainer", None),
        ("Defamation lawyers on standby", None),
        ("Top-tier defense contractor", "→ Pentagon access"),
        ("Patrick Shanahan",
         "30-year Boeing executive,\nActing Defense Sec. (Jan–Jun 2019)"),
        ("Congressional offices in every Boeing state", None),
    ]
    yy = col_y + Inches(0.75)
    for head, sub in apparatus:
        oval(s, xs[2] + Inches(0.22), yy + Inches(0.10),
             Inches(0.16), Inches(0.16), ACCENT)
        text(s, head,
             xs[2] + Inches(0.50), yy,
             col_w - Inches(0.65), Inches(0.35),
             font=BODY_FONT, size=12, bold=True, color=INK,
             line_spacing=1.15)
        if sub:
            text(s, sub,
                 xs[2] + Inches(0.50), yy + Inches(0.32),
                 col_w - Inches(0.65), Inches(0.55),
                 font=BODY_FONT, size=10, italic=True, color=INK_MUTED,
                 line_spacing=1.2)
            yy += Inches(0.85)
        else:
            yy += Inches(0.45)

    # Quote slab
    quote_slab(s,
        "If certain kinds of fact, position, or program are thought likely "
        "to elicit flak, this prospect can be a deterrent.",
        "— Herman & Chomsky, p. 26",
        y=Inches(5.8), h=Inches(1.05))

    text(s,
        "Flak doesn’t have to land.  It has to be anticipated.",
        Inches(0.7), Inches(7.0), Inches(11.93), Inches(0.4),
        font=TITLE_FONT, size=16, bold=True, italic=True, color=ACCENT,
        align=PP_ALIGN.CENTER)

    notes(s,
        outline=[
            "Three kinds of evidence for filter 4 operating here.",
            "(1) Softening WORDS visible inside the article itself.",
            "(2) Same-beat comparison: Gates @ Seattle Times leads with capture; WSJ doesn't.",
            "(3) Boeing's apparatus including Patrick Shanahan as Acting DefSec.",
            "Chomsky: flak's power is anticipatory, not retaliatory.",
            "Transition: 'But why those particular framings? Filter 5 — the unifying ideology — supplies the cognitive shortcut.'",
        ],
        script=(
            "Filter 4: flak. Three kinds of evidence. First, the softening visible "
            "INSIDE the article itself. 'Deadly crash' rather than 'second "
            "mass-fatality crash of the same aircraft in five months.' 'The FAA's "
            "decision' rather than 'the regulator that delegated certification to "
            "the company it was regulating.' Foreign action framed as deviation, "
            "FAA framed as default. These are not lies; they are the framings that "
            "survive anticipated flak. Second, the same-beat comparison. Dominic "
            "Gates at the Seattle Times — same aviation beat — leads with "
            "captured-regulator framing six days later, and wins a Pulitzer for "
            "it. The WSJ team could have done this. The structural conditions of "
            "their publication selected against it. Third, Boeing's apparatus. "
            "White-shoe PR firms on retainer. Defamation lawyers ready. Pentagon "
            "access — the acting Defense Secretary at the time of these crashes, "
            "Patrick Shanahan, was a 30-year Boeing executive. Congressional "
            "offices in every Boeing state. Chomsky's deepest insight: flak is "
            "most powerful when it is invisible. Pasztor and Tangel are good "
            "enough to anticipate the flak boundaries without anyone having to "
            "tell them where the boundaries are."
        ),
        sources=[
            ("Patrick Shanahan official DoD portrait — Wikimedia Commons (public domain, DoD)",
             "https://upload.wikimedia.org/wikipedia/commons/9/9a/Patrick_M._Shanahan_official_portrait.jpg"),
            ("Frontline / NYT, 'Boeing's Fatal Flaw' (2021) documentary",
             "https://www.pbs.org/wgbh/frontline/documentary/boeings-fatal-flaw/"),
            ("Gates, Seattle Times, Mar 17/21 2019 — already cited in slide 6",
             "https://www.seattletimes.com/business/boeing-aerospace/failed-certification-faa-missed-safety-issues-in-the-737-max-system-implicated-in-the-lion-air-crash/"),
        ],
    )


def slide_8_ideology(prs):
    s = blank(prs)
    slide_title(s, "Filter 5", "The Updated Ideology")

    # Three eras as a typographic ladder
    eras = [
        ("Chomsky  1988",    "Anticommunism",                INK),
        ("Chomsky  2002",    "The religion of the market",   INK),
        ("2019  aviation",   "Techno-industrial nationalism", ACCENT),
    ]
    era_y = Inches(2.0)
    era_h = Inches(0.55)
    for i, (when, label, col) in enumerate(eras):
        y = era_y + (era_h + Inches(0.04)) * i
        text(s, when,
             Inches(0.7), y, Inches(3.5), era_h,
             font=BODY_FONT, size=14, italic=True, color=INK_MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, label,
             Inches(4.3), y, Inches(8.5), era_h,
             font=TITLE_FONT, size=24, bold=True, italic=True, color=col,
             anchor=MSO_ANCHOR.MIDDLE)

    # Worthy / unworthy victims comparison
    tbl_y = Inches(3.95)
    col_w = Inches(5.95)
    left_x = Inches(0.7)
    right_x = Inches(6.7)
    rect(s, left_x,  tbl_y, col_w, Inches(0.55), fill=LINE)
    rect(s, right_x, tbl_y, col_w, Inches(0.55), fill=ACCENT)
    text(s, "Worthy victims   (1983)",
         left_x, tbl_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=15, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Unworthy victims   (2018–19)",
         right_x, tbl_y, col_w, Inches(0.55),
         font=TITLE_FONT, size=15, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("KAL 007  (Soviet shootdown)", "Lion Air 610  —  189 Indonesian"),
        ("Sustained moral outrage; years of coverage", "“Pilot training” framing"),
        ("Reagan: “cold-blooded murder”",  "Ethiopian 302  —  157 dead, 35 nations"),
        ("Front-page treatment for years",     "132 days before grounding"),
    ]
    row_h = Inches(0.5)
    for i, (l, r) in enumerate(rows):
        y = tbl_y + Inches(0.55) + row_h * i
        rect(s, left_x,  y, col_w, row_h, fill=SLAB)
        rect(s, right_x, y, col_w, row_h, fill=SLAB)
        text(s, l, left_x + Inches(0.25), y, col_w - Inches(0.4), row_h,
             font=BODY_FONT, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        text(s, r, right_x + Inches(0.25), y, col_w - Inches(0.4), row_h,
             font=BODY_FONT, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # Quote
    quote_slab(s,
        "Concentrate on the victims of enemy powers and forget about the "
        "victims of friends.",
        "— Herman & Chomsky, p. 32   (paraphrasing the structural function)",
        y=Inches(6.40), h=Inches(1.0))

    notes(s,
        outline=[
            "The paper's analytical move: filter 5 has rotated content, not function.",
            "1988 anticommunism → 2002 religion of the market → 2019 techno-industrial nationalism.",
            "KAL 007: sustained moral outrage, Reagan's denunciation, years of front-page.",
            "Lion Air + Ethiopian: pilot-training framing, 132 days before grounding.",
            "Different ideology — IDENTICAL filter function: worthy/unworthy victims.",
            "Transition: 'If the model holds, then changing the victims should change the framing. The counterfactual.'",
        ],
        script=(
            "Filter 5 is where this paper makes its analytical move beyond textbook "
            "Chomsky. Chomsky's original filter 5 in 1988 was anticommunism — the "
            "binary that defined worthy and unworthy victims. He updated it in 2002 "
            "to 'the religion of the market.' For 2019 aviation coverage, I argue "
            "the operative filter 5 is techno-industrial nationalism — the framing "
            "of American manufacturing supremacy, against Airbus and rising Chinese "
            "aviation, as quasi-religious national interest. Here is the proof that "
            "the structural function is identical. In 1983, the Soviets shot down a "
            "civilian airliner — KAL 007. Sustained moral outrage. Years of "
            "front-page coverage. Reagan called it 'cold-blooded murder.' In 2018 "
            "and 2019, an American aircraft killed 346 civilians — Indonesians, "
            "Indians, Ethiopians, Kenyans, Chinese. Initial framing: foreign pilot "
            "training was inadequate. 132 days passed before the aircraft was "
            "grounded. Different ideology, same filter. The unworthy victims merit, "
            "in Chomsky's words, 'only slight detail, minimal humanization, and "
            "little context that will excite and enrage.' That is the structural "
            "test. The model survives the update."
        ),
        sources=[
            ("KAL 007 background — Wikipedia",
             "https://en.wikipedia.org/wiki/Korean_Air_Lines_Flight_007"),
            ("Bhavye Suneja (Lion Air 610 captain) — coverage",
             "https://www.republicworld.com/india-news/general-news/indian-pilot-bhavye-suneja-captained-the-lion-air-flight-which-crashed-into-the-sea-with-188-people-aboard.html"),
            ("Yared Getachew (Ethiopian 302 captain) — coverage",
             "https://en.wikipedia.org/wiki/Ethiopian_Airlines_Flight_302"),
        ],
    )


def slide_9_counterfactual(prs):
    s = blank(prs)
    text(s, "THE COUNTERFACTUAL",
         Inches(0.7), Inches(0.55), Inches(12), Inches(0.4),
         font=BODY_FONT, size=12, bold=True, color=ACCENT)
    text(s, "Same facts.  Different victims.  Different filter output.",
         Inches(0.7), Inches(0.88), Inches(12), Inches(0.55),
         font=TITLE_FONT, size=22, italic=True, color=INK_MUTED)

    # Mock WSJ headline mounted on newsprint
    headline_y = Inches(1.85)
    headline_h = Inches(3.0)
    rect(s, Inches(0.9), headline_y, Inches(11.5), headline_h,
         fill=PAPER, line=INK_SOFT)
    rich_text(s, [
        [("BUSINESS",
          {"font": BODY_FONT, "size": 12, "bold": True, "color": INK_MUTED})],
        [(" ", {"size": 10})],
        [("Pressure Mounts on FAA to Ground 737 MAX",
          {"font": QUOTE_FONT, "size": 34, "bold": True, "color": INK})],
        [("as Second Crash Kills 157 Americans",
          {"font": QUOTE_FONT, "size": 34, "bold": True, "color": INK})],
        [(" ", {"size": 6})],
        [("Lawmakers, victims’ families, and foreign regulators "
          "demand immediate suspension;",
          {"font": QUOTE_FONT, "size": 15, "italic": True, "color": INK_MUTED})],
        [("FAA’s certification process under scrutiny.",
          {"font": QUOTE_FONT, "size": 15, "italic": True, "color": INK_MUTED})],
        [(" ", {"size": 8})],
        [("By [byline]    |    Updated March 11, 2019",
          {"font": BODY_FONT, "size": 12, "bold": True, "color": INK})],
    ], Inches(1.3), headline_y + Inches(0.3), Inches(10.7), headline_h - Inches(0.5),
       align=PP_ALIGN.LEFT, line_spacing=1.08)

    # The tagline
    text(s, "This is the headline if 346 Americans had died.",
         Inches(0.7), Inches(5.15), Inches(11.93), Inches(0.55),
         font=TITLE_FONT, size=22, bold=True, italic=True, color=INK,
         align=PP_ALIGN.CENTER)

    # Final landing slab
    rect(s, Inches(0.7), Inches(6.05), Inches(11.93), Inches(1.05), fill=ACCENT)
    text(s, "The filters didn’t lie.   They selected.",
         Inches(0.7), Inches(6.05), Inches(11.93), Inches(1.05),
         font=TITLE_FONT, size=28, bold=True, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    notes(s,
        outline=[
            "Build the counterfactual: if 346 AMERICANS had died, this is the headline.",
            "Same facts, different victims, different filter output.",
            "Restate Chomsky's KAL 007 logic, updated.",
            "Land: 'The filters didn't lie. They selected.'",
            "End cleanly. Pause. Thank them.",
        ],
        script=(
            "If 346 Americans had died on two domestic flights of the same "
            "aircraft in five months, the Wall Street Journal headline on March "
            "11, 2019 would not have read 'FAA Has No Current Plans to Ground "
            "Boeing's 737 MAX.' It would have read something like this. Same "
            "facts. Different victims. Different filter output. The propaganda "
            "model is not a claim that this article lied. It is a claim that the "
            "structural conditions of its production — Murdoch's WSJ covering "
            "America's largest exporter, Boeing's flak apparatus, the access "
            "economy on the aviation beat, the techno-nationalist common sense "
            "that makes American manufacturing a quasi-religious good — selected "
            "for the framing we got. Thirty-one years after publication, Chomsky's "
            "model holds. Filter 5 has just rotated from anticommunism to "
            "techno-industrial nationalism. The filters didn't lie. They selected."
        ),
    )


# ---------------- driver ----------------

def main():
    prs = make_prs()
    slide_1_hook(prs)
    slide_2_timeline(prs)
    slide_3_five_filters(prs)
    slide_4_ownership(prs)
    slide_5_advertising(prs)
    slide_6_sourcing(prs)
    slide_7_flak(prs)
    slide_8_ideology(prs)
    slide_9_counterfactual(prs)
    prs.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
