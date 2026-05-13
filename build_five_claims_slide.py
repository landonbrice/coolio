"""Single-slide PowerPoint artifact: 'Five claims hold up the paper.'

Styled to match the author's reference slide (Aptos sans, red accent,
thin blue rule, numbered-claim grid, italic closing quote).  Open the
output .pptx in PowerPoint -- 'Aptos' will resolve to the modern
default sans on Microsoft 365.  On systems without Aptos, PowerPoint
substitutes Calibri.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
OUT  = ROOT / "Boeing-WSJ-Five-Claims.pptx"

SANS  = "Aptos"
SANS_DISPLAY = "Aptos Display"
INK   = RGBColor(0x11, 0x11, 0x11)
RED   = RGBColor(0xC0, 0x00, 0x00)
GREY  = RGBColor(0x4F, 0x4F, 0x4F)
MUTE  = RGBColor(0x6B, 0x6B, 0x6B)
RULE  = RGBColor(0x9D, 0xC3, 0xE6)
SOFT  = RGBColor(0xBF, 0xD7, 0xEE)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def textbox(x, y, w, h, runs, anchor="top"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]

    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        for j, (txt, opts) in enumerate(line):
            r = p.add_run() if (i or j) else p.add_run()
            # first run on first paragraph is auto-created; reuse it
            if i == 0 and j == 0:
                r = p.runs[0] if p.runs else p.add_run()
            r.text = txt
            f = r.font
            f.name = opts.get("font", SANS)
            f.size = Pt(opts.get("size", 18))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", INK)
        if "space_after" in line[0][1]:
            p.space_after = Pt(line[0][1]["space_after"])
        if "line_spacing" in line[0][1]:
            p.line_spacing = line[0][1]["line_spacing"]
    return tb


def hrule(x, y, w, weight=1.75, color=RULE):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


PAD_L = Inches(0.55)
PAD_R = Inches(0.55)
CONTENT_W = prs.slide_width - PAD_L - PAD_R
y = Inches(0.40)

# Title
textbox(PAD_L, y, CONTENT_W, Inches(0.78),
        [[("Five claims hold up the paper", {"font": SANS_DISPLAY, "size": 34, "bold": True, "color": INK, "line_spacing": 1.05}),
          (".", {"font": SANS_DISPLAY, "size": 34, "bold": True, "color": RED})]])
y += Inches(0.62)

# Subtitle
textbox(PAD_L, y, CONTENT_W, Inches(0.30),
        [[("descriptive, structural, analytical, update, proof.",
           {"size": 14, "italic": True, "color": GREY})]])
y += Inches(0.34)

# Rule
hrule(PAD_L, y, CONTENT_W, weight=1.75, color=RULE)
y += Inches(0.14)

# Lead paragraph (3 lines at 13pt)
textbox(PAD_L, y, CONTENT_W, Inches(0.95),
        [[("The paper makes five claims. Each is independently defensible from the article text and from public structural facts about the WSJ, Boeing, and the aviation reporting beat. Together they make the propaganda model a sufficient explanation of the March 11, 2019 frame.",
           {"size": 13, "color": INK, "line_spacing": 1.30})]])
y += Inches(1.02)

# Claims grid — variable row heights
LABEL_W = Inches(2.00)
BODY_X  = PAD_L + LABEL_W + Inches(0.22)
BODY_W  = prs.slide_width - BODY_X - PAD_R

claims = [
    ("I — Descriptive",
     "Every one of Chomsky's five filters operates on the WSJ article of March 11, 2019. ",
     "The operation can be demonstrated from the article text and from public structural facts about the WSJ, Boeing, and the aviation reporting beat.",
     0.72),
    ("II — Structural",
     "The filters do not operate independently. ",
     "They reinforce one another such that no individual filter is decisive. The joint operation produces the cleansed residue.",
     0.55),
    ("III — Analytical distinction",
     "The journalists are not the explanatory variable. ",
     "The filters operate even when reporters are competent, sincere, and individually critical — which Pasztor and Tangel demonstrably are. This is what separates a sophisticated application of Chomsky from a naive one.",
     0.92),
    ("IV — Model update",
     "Filter 5 has rotated from anticommunism (1988) to techno-industrial nationalism (2019) without changing its structural function. ",
     "The worthy / unworthy victims distinction operates identically. This is the paper's analytical contribution.",
     0.75),
    ("V — Proof",
     "Same facts, different victims, different output. ",
     "If 346 Americans had died, the same filters would have produced different output, because the worthy-victim category determines what the cleansed residue looks like. The model is confirmed by the structural pattern of which victims get which kinds of coverage, not refuted by any single article.",
     0.95),
]

for label, opener, rest, h in claims:
    row_h = Inches(h)
    textbox(PAD_L, y, LABEL_W, row_h,
            [[(label.upper(),
               {"size": 10, "bold": True, "color": RED})]],
            anchor="top")
    textbox(BODY_X, y, BODY_W, row_h,
            [[(opener, {"size": 13, "bold": True, "color": INK, "line_spacing": 1.28}),
              (rest,   {"size": 13, "color": INK, "line_spacing": 1.28})]],
            anchor="top")
    y += row_h + Inches(0.06)

# Closing quote — anchored above the bottom margin
QY = Inches(6.95)
hrule(PAD_L, QY, CONTENT_W, weight=0.75, color=SOFT)
textbox(PAD_L, QY + Inches(0.06), CONTENT_W, Inches(0.40),
        [[("“The constraints are so powerful, and are built into the system in such a fundamental way, that alternative bases of news choices are hardly imaginable.”  ",
           {"size": 11, "italic": True, "color": GREY}),
          ("— Herman & Chomsky, Manufacturing Consent, p. 2",
           {"size": 11, "color": MUTE})]])

prs.save(OUT)
print(f"wrote {OUT.name} ({OUT.stat().st_size:,} bytes)")
